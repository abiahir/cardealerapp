"""PowerPoint presentation generation logic."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .models import DealerDetails, VehicleListing


class PresentationBuilder:
    """Builds PowerPoint presentations from vehicle listings."""

    @staticmethod
    def build(vehicle: VehicleListing, output_path: Path) -> None:
        """Build and save a presentation for the given vehicle."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        PresentationBuilder._add_title_block(slide, vehicle)
        PresentationBuilder._add_specification_table(slide, vehicle)
        PresentationBuilder._add_contact_block(slide, vehicle.dealer)

        prs.save(output_path)

    @staticmethod
    def _add_title_block(slide, vehicle: VehicleListing) -> None:
        """Add the title block with dealer name, price, and vehicle title."""
        # Dealer name - top left, large and bold
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6.5), Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = vehicle.dealer.name
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

        # Price pill - top right, red with white text
        pill = slide.shapes.add_shape(
            1,  # MSO_SHAPE_RECTANGLE
            Inches(7.5),
            Inches(0.3),
            Inches(2.0),
            Inches(0.8),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(0xF7, 0x00, 0x00)
        pill.line.color.rgb = RGBColor(0xF7, 0x00, 0x00)  # Match fill for seamless look
        pill.line.width = Pt(0)  # No border

        pill_text_frame = pill.text_frame
        pill_text_frame.vertical_anchor = 1  # Middle vertical alignment
        pill_text = pill_text_frame.paragraphs[0]
        pill_text.text = vehicle.price
        pill_text.font.size = Pt(32)
        pill_text.font.bold = True
        pill_text.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        pill_text.alignment = PP_ALIGN.CENTER

        # Vehicle title - below dealer name
        subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.7))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = vehicle.title
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

    @staticmethod
    def _add_specification_table(slide, vehicle: VehicleListing) -> None:
        """Add the specification table with vehicle details."""
        left, top = Inches(0.5), Inches(2.0)
        width, height = Inches(9), Inches(6.5)
        table_shape = slide.shapes.add_table(rows=10, cols=2, left=left, top=top, width=width, height=height)
        table = table_shape.table

        labels = [
            "Registration",
            "Year",
            "Gearbox",
            "Engine Size",
            "Fuel Type",
            "Owners",
            "Mileage",
            "ULEZ",
            "MOT Expiry",
            "Specs",
        ]
        values = [
            vehicle.registration,
            vehicle.year,
            vehicle.gearbox,
            vehicle.engine_size,
            vehicle.fuel_type,
            vehicle.owners,
            vehicle.mileage,
            vehicle.ulez,
            vehicle.mot_expiry,
            ", ".join(vehicle.specs) if vehicle.specs else "",
        ]

        header_fill = RGBColor(0x00, 0x38, 0x64)
        value_fill = RGBColor(0xEC, 0xEC, 0xEC)

        for row_idx, (label, value) in enumerate(zip(labels, values)):
            label_cell = table.cell(row_idx, 0)
            value_cell = table.cell(row_idx, 1)

            # Set cell text
            label_cell.text = label
            value_cell.text = value

            # Configure label cell (blue background, white text)
            label_cell.fill.solid()
            label_cell.fill.fore_color.rgb = header_fill
            label_para = label_cell.text_frame.paragraphs[0]
            label_para.font.size = Pt(20)
            label_para.font.bold = True
            label_para.font.name = "Calibri"
            label_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            label_para.alignment = PP_ALIGN.LEFT
            label_cell.text_frame.vertical_anchor = 1  # Middle vertical alignment
            label_cell.margin_left = Inches(0.1)
            label_cell.margin_right = Inches(0.1)
            label_cell.margin_top = Inches(0.05)
            label_cell.margin_bottom = Inches(0.05)

            # Configure value cell (gray background, black text)
            value_cell.fill.solid()
            value_cell.fill.fore_color.rgb = value_fill
            value_para = value_cell.text_frame.paragraphs[0]
            value_para.font.size = Pt(20)
            value_para.font.name = "Calibri"
            value_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            value_para.alignment = PP_ALIGN.LEFT
            value_cell.text_frame.vertical_anchor = 1  # Middle vertical alignment
            value_cell.margin_left = Inches(0.1)
            value_cell.margin_right = Inches(0.1)
            value_cell.margin_top = Inches(0.05)
            value_cell.margin_bottom = Inches(0.05)

        # Set column widths
        table.columns[0].width = Inches(3.2)
        table.columns[1].width = Inches(5.8)

    @staticmethod
    def _add_contact_block(slide, dealer: DealerDetails) -> None:
        """Add the contact information block at the bottom."""
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.7), Inches(9), Inches(0.6))
        p = contact_box.text_frame.paragraphs[0]
        p.text = f"Call: {dealer.phone}  |  Email: {dealer.email}  |  Web: {dealer.website}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)

