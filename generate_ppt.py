"""Generate an editable vehicle presentation from structured data.

The script reads vehicle details from JSON or CLI flags and builds a
PowerPoint slide that mirrors a showroom summary. All text boxes remain
fully editable so that last-minute changes can be made directly in the
.PPTX file.
"""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Supported categorical options exposed via the CLI.
FUEL_TYPES = ["Petrol", "Diesel", "Hybrid", "Electric"]
GEARBOX_TYPES = ["Automatic", "Manual"]
ULEZ_CHOICES = ["Yes", "No", "Unknown"]


@dataclass
class DealerDetails:
    name: str = "Your Dealership"
    phone: str = "0000 000 0000"
    email: str = "sales@example.com"
    website: str = "www.example.com"

    @staticmethod
    def from_dict(data: Optional[Dict[str, str]]) -> "DealerDetails":
        if not data:
            return DealerDetails()
        return DealerDetails(
            name=data.get("name", "Your Dealership"),
            phone=data.get("phone", "0000 000 0000"),
            email=data.get("email", "sales@example.com"),
            website=data.get("website", "www.example.com"),
        )


@dataclass
class VehicleListing:
    title: str
    price: str
    registration: str
    year: str
    gearbox: str
    engine_size: str
    fuel_type: str
    mileage: str
    ulez: str
    mot_expiry: str
    owners: str = "Unknown"
    specs: List[str] = field(default_factory=list)
    dealer: DealerDetails = field(default_factory=DealerDetails)

    @staticmethod
    def from_dict(data: Dict[str, object]) -> "VehicleListing":
        def _choice(value: str, allowed: List[str], label: str) -> str:
            if value not in allowed:
                choices = ", ".join(allowed)
                raise ValueError(f"{label} must be one of: {choices}")
            return value

        return VehicleListing(
            title=str(data.get("title", "Vehicle Title")),
            price=str(data.get("price", "Price on enquiry")),
            registration=str(data.get("registration", "Registration")),
            year=str(data.get("year", "Year")),
            gearbox=_choice(str(data.get("gearbox", "Automatic")), GEARBOX_TYPES, "Gearbox"),
            engine_size=str(data.get("engine_size", "2.0 L")),
            fuel_type=_choice(str(data.get("fuel_type", "Diesel")), FUEL_TYPES, "Fuel type"),
            mileage=str(data.get("mileage", "0")),
            ulez=_choice(str(data.get("ulez", "Unknown")), ULEZ_CHOICES, "ULEZ"),
            mot_expiry=str(data.get("mot_expiry", "Unknown")),
            owners=str(data.get("owners", "Unknown")),
            specs=[str(item) for item in data.get("specs", [])],
            dealer=DealerDetails.from_dict(data.get("dealer")),
        )

    def update_from_args(self, args: argparse.Namespace) -> None:
        """Override fields with CLI arguments when provided."""
        for field_name in (
            "title",
            "price",
            "registration",
            "year",
            "gearbox",
            "engine_size",
            "fuel_type",
            "mileage",
            "ulez",
            "mot_expiry",
            "owners",
        ):
            value = getattr(args, field_name)
            if value is not None:
                setattr(self, field_name, value)

        if args.specs:
            self.specs = args.specs

        self.dealer = DealerDetails(
            name=args.dealer_name or self.dealer.name,
            phone=args.dealer_phone or self.dealer.phone,
            email=args.dealer_email or self.dealer.email,
            website=args.dealer_website or self.dealer.website,
        )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate an editable vehicle PPTX.")
    parser.add_argument("--input", type=Path, help="Path to vehicle JSON data.")
    parser.add_argument("--output", type=Path, default=Path("vehicle.pptx"), help="PPTX output file.")

    parser.add_argument("--title", help="Vehicle display title.")
    parser.add_argument("--price", help="Price string, e.g. '£12,000'.")
    parser.add_argument("--registration", help="Registration number.")
    parser.add_argument("--year", help="Year display value.")
    parser.add_argument("--gearbox", choices=GEARBOX_TYPES, help="Gearbox option.")
    parser.add_argument("--engine-size", dest="engine_size", help="Engine size label.")
    parser.add_argument("--fuel-type", dest="fuel_type", choices=FUEL_TYPES, help="Fuel type option.")
    parser.add_argument("--mileage", help="Mileage display value.")
    parser.add_argument("--ulez", choices=ULEZ_CHOICES, help="ULEZ status.")
    parser.add_argument("--mot-expiry", dest="mot_expiry", help="MOT expiry display value.")
    parser.add_argument("--owners", help="Number of owners or description.")
    parser.add_argument("--specs", nargs="*", help="List of spec bullet points.")

    parser.add_argument("--dealer-name", dest="dealer_name", help="Dealership name.")
    parser.add_argument("--dealer-phone", dest="dealer_phone", help="Dealership phone number.")
    parser.add_argument("--dealer-email", dest="dealer_email", help="Dealership email address.")
    parser.add_argument("--dealer-website", dest="dealer_website", help="Dealership website.")

    return parser.parse_args()


def load_vehicle(path: Optional[Path]) -> VehicleListing:
    if path is None:
        return VehicleListing.from_dict({})
    data = json.loads(path.read_text())
    if not isinstance(data, dict):
        raise ValueError("Input JSON must describe a single vehicle object.")
    return VehicleListing.from_dict(data)


def add_title_block(slide, vehicle: VehicleListing) -> None:
    # Dealer name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = vehicle.dealer.name
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Price pill
    pill = slide.shapes.add_shape(
        1,  # MSO_SHAPE_RECTANGLE
        Inches(7.5),
        Inches(0.35),
        Inches(2.0),
        Inches(0.9),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(0xF7, 0x00, 0x00)
    pill.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)

    pill_text = pill.text_frame.paragraphs[0]
    pill_text.text = vehicle.price
    pill_text.font.size = Pt(28)
    pill_text.font.bold = True
    pill_text.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    pill_text.alignment = PP_ALIGN.CENTER

    # Vehicle title
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.8))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = vehicle.title
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)


def add_specification_table(slide, vehicle: VehicleListing) -> None:
    left, top = Inches(0.5), Inches(2.2)
    width, height = Inches(9), Inches(7)
    table_shape = slide.shapes.add_table(rows=10, cols=2, left=left, top=top, width=width, height=height)
    table = table_shape.table

    labels = [
        "Registration",
        "Year",
        "Gearbox",
        "Engine Size",
        "Fuel Type",
        "Owners",
        "Mileage",
        "ULEZ",
        "MOT Expiry",
        "Specs",
    ]
    values = [
        vehicle.registration,
        vehicle.year,
        vehicle.gearbox,
        vehicle.engine_size,
        vehicle.fuel_type,
        vehicle.owners,
        vehicle.mileage,
        vehicle.ulez,
        vehicle.mot_expiry,
        ", ".join(vehicle.specs) if vehicle.specs else "",
    ]

    header_fill = RGBColor(0x00, 0x38, 0x64)
    for row_idx, (label, value) in enumerate(zip(labels, values)):
        label_cell = table.cell(row_idx, 0)
        value_cell = table.cell(row_idx, 1)

        label_cell.text = label
        value_cell.text = value

        for cell in (label_cell, value_cell):
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = "Calibri"

        label_cell.fill.solid()
        label_cell.fill.fore_color.rgb = header_fill
        label_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        value_cell.fill.solid()
        value_cell.fill.fore_color.rgb = RGBColor(0xEC, 0xEC, 0xEC)
        value_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(6)


def add_contact_block(slide, dealer: DealerDetails) -> None:
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(9.3), Inches(9), Inches(0.8))
    p = contact_box.text_frame.paragraphs[0]
    p.text = f"Call: {dealer.phone}  |  Email: {dealer.email}  |  Web: {dealer.website}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER


def build_presentation(vehicle: VehicleListing, output_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    add_title_block(slide, vehicle)
    add_specification_table(slide, vehicle)
    add_contact_block(slide, vehicle.dealer)

    prs.save(output_path)


def main() -> None:
    args = parse_args()
    vehicle = load_vehicle(args.input)
    vehicle.update_from_args(args)
    build_presentation(vehicle, args.output)
    print(f"Saved editable PPTX to {args.output}")


if __name__ == "__main__":
    main()
